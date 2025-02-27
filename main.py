import os
import subprocess
import sys
from dotenv import load_dotenv
from langchain_google_genai import GoogleGenerativeAI
from langchain.prompts import PromptTemplate
from PyPDF2 import PdfReader
from pptx import Presentation
import shutil
import textwrap
import pyfiglet

# Load environment variables from .env file
load_dotenv()
gemini_api_key = os.getenv("GEMINI_API_KEY")
if not gemini_api_key:
    raise ValueError("Gemini API key is not set in the environment variables.")

# Initialize the Gemini model
llm = GoogleGenerativeAI(model="gemini-1.5-pro")

# Define the prompt templates
summary_prompt_template = """
Given the following text from a study document, extract the key points and create a concise summary and bullet points for an examination reviewer about this.
Text:
{text}
Please format your response exactly as follows:
Summary:
[A brief overview of the main topics and key concepts]
Bullet Points:
• [Key point 1]
• [Key point 2]
• [Key point 3]
...and so on
Do not use any other formatting, just plain bullet points and text.
"""

qa_prompt_template = """
Given the following text and a question, provide a comprehensive and accurate answer based solely on the information in the text.

Text:
{text}

Question:
{question}

Answer the question directly and concisely, using only information provided in the text. If the answer cannot be determined from the text, state that clearly.
"""

# Initialize the PromptTemplates
summary_prompt = PromptTemplate(input_variables=["text"], template=summary_prompt_template)
qa_prompt = PromptTemplate(input_variables=["text", "question"], template=qa_prompt_template)

def extract_text_from_pdf(file_path):
    """Extract text from each page of the PDF."""
    pdf_text = []
    try:
        with open(file_path, "rb") as file:
            reader = PdfReader(file)
            for page in reader.pages:
                text = page.extract_text()
                if text:  # Check if text extraction was successful
                    pdf_text.append(text)
    except Exception as e:
        raise IOError(f"Error reading the PDF file: {e}")
    return " ".join(pdf_text)  # Join the list into a single string

def convert_ppt_to_pptx(ppt_path):
    """Convert .ppt to .pptx using available tools."""
    # create tmp file with temp_ prefix
    temp_dir = os.path.dirname(ppt_path)
    temp_filename = f"temp_{os.path.basename(ppt_path)}x"
    pptx_path = os.path.join(temp_dir, temp_filename)
    manual_conversion_msg = """
No conversion tools found. You have these options:
1. Manually convert the .ppt file to .pptx:
   - Open the file in Microsoft PowerPoint, LibreOffice, or WPS Office
   - Save As -> Select .pptx format
   - Try again with the converted file
2. Install one of these tools:
   - LibreOffice (Recommended): Install from your system's package manager
   - unoserver: pip install unoserver (requires LibreOffice)
   - unoconv: Available in package manager (requires LibreOffice)
"""
    try:
        if sys.platform == 'win32':
            try:
                print("Attempting conversion using Microsoft PowerPoint...")
                import win32com.client
                powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                powerpoint.Visible = False
                deck = powerpoint.Presentations.Open(ppt_path)
                deck.SaveAs(pptx_path, 24)  # 24 = .pptx format
                deck.Close()
                powerpoint.Quit()
                print("Successfully converted using Microsoft PowerPoint")
                return pptx_path
            except Exception as e:
                raise IOError(f"PowerPoint automation failed: {e}\n{manual_conversion_msg}")
        else:
            # conversion tools in order of preference
            converters = [
                {
                    'name': 'unoserver',
                    'cmd': ['unoconvert', ppt_path, pptx_path],
                    'install': 'pip install unoserver (requires LibreOffice)'
                },
                {
                    'name': 'unoconv',
                    'cmd': ['unoconv', '-f', 'pptx', '-o', pptx_path, ppt_path],
                    'install': 'Install from package manager (requires LibreOffice)'
                },
                {
                    'name': 'soffice',
                    'cmd': ['soffice', '--headless', '--convert-to', 'pptx', '--outdir', temp_dir, ppt_path],
                    'install': 'Install LibreOffice'
                }
            ]
            errors = []
            for converter in converters:
                if shutil.which(converter['cmd'][0]):  # Check if command exists
                    try:
                        print(f"Attempting conversion using {converter['name']}...")
                        result = subprocess.run(
                            converter['cmd'],
                            capture_output=True,
                            text=True
                        )
                        if result.returncode == 0 and os.path.exists(pptx_path):
                            print(f"Successfully converted using {converter['name']}")
                            return pptx_path
                        else:
                            errors.append(f"{converter['name']}: Command failed with return code {result.returncode}")
                    except Exception as e:
                        errors.append(f"{converter['name']}: {str(e)}")
                        print(f"Failed to convert using {converter['name']}")
                        continue
                else:
                    print(f"{converter['name']} not found, trying next method...")
            # no converter worked
            error_msg = "\n".join([
                "Failed to convert .ppt to .pptx automatically.",
                manual_conversion_msg,
                "\nTechnical details:",
                *errors
            ])
            raise IOError(error_msg)
    except Exception as e:
        raise e

def extract_text_from_pptx(file_path):
    """Extract text from PowerPoint presentation."""
    pptx_text = []
    file_ext = os.path.splitext(file_path)[1].lower()
    temp_file = None
    try:
        if file_ext == '.ppt':
            print("\nNote: .ppt file detected, attempting automatic conversion...")
            temp_file = convert_ppt_to_pptx(file_path)
            file_path = temp_file
        print("\nExtracting text from presentation...")
        prs = Presentation(file_path)
        # Process slides
        for slide_number, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                try:
                    text = str(shape.text).strip()  # type: ignore
                    if text:
                        slide_text.append(text)
                except:
                    continue
            if slide_text:
                pptx_text.append(f"[Slide {slide_number}] {' '.join(slide_text)}")
    except Exception as e:
        error_msg = "Error reading PowerPoint file"
        if file_ext == '.ppt':
            error_msg = "Error reading .ppt file"
        raise IOError(f"{error_msg}: {e}")
    finally:
        # rm tmp file if it exists and prefix with 'temp_'
        if temp_file and os.path.exists(temp_file) and os.path.basename(temp_file).startswith('temp_'):
            try:
                os.remove(temp_file)
            except:
                pass
    return "\n\n".join(pptx_text)

def extract_text_from_file(file_path):
    """Extract text based on file type."""
    file_ext = os.path.splitext(file_path)[1].lower()
    if file_ext == '.pdf':
        return extract_text_from_pdf(file_path)
    elif file_ext in ['.pptx', '.ppt']:
        return extract_text_from_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file format: {file_ext}")

def summarize_text_to_bullets(text):
    """Summarize the extracted text into bullet points using the LLM."""
    chain = summary_prompt | llm
    response = chain.invoke({"text": text})
    return response

def answer_question(text, question):
    """Answer a question based on the document content."""
    chain = qa_prompt | llm
    response = chain.invoke({"text": text, "question": question})
    return response

def process_file(file_path):
    """Extract text from file and return it."""
    return extract_text_from_file(file_path)

def get_file_path():
    """Get file path from command line argument or user input."""
    def is_valid_file(path):
        """Check if the path exists and has a supported extension."""
        if not os.path.exists(path):
            print(f"Error: File '{path}' does not exist")
            return False
        valid_extensions = {
            '.pdf': 'PDF Document',
            '.pptx': 'PowerPoint Document',
            '.ppt': 'Legacy PowerPoint Document (Limited Support)'
        }
        file_ext = os.path.splitext(path)[1].lower()
        if file_ext not in valid_extensions:
            print(f"Error: Unsupported file format. Supported formats are:")
            for ext, desc in valid_extensions.items():
                print(f"  {ext} - {desc}")
            return False
        return True
    
    # Check cli
    if len(sys.argv) > 1:
        file_path = sys.argv[1]
        if is_valid_file(file_path):
            return file_path
        sys.exit(1)  # if argument is invalid
    
    # prompt for input
    while True:
        file_path = input("\nPlease enter the path to your file (.pdf, .pptx, or .ppt): ").strip()
        if file_path.lower() == 'exit':
            print("Exiting program...")
            sys.exit(0)
        if is_valid_file(file_path):
            return file_path
        print("\nTry again or type 'exit' to quit")

def display_menu():
    """Display the menu"""
    print("\n"+pyfiglet.figlet_format("REVIEWER", "slant"))
    print("[1] Create Summary and Bullet Points")
    print("[2] Ask Questions About the Document")
    print("[3] Exit")
    while True:
        try:
            choice = int(input("\nEnter your choice (1-3): "))
            if 1 <= choice <= 3:
                return choice
            print("Invalid choice. Please enter a number between 1 and 3.")
        except ValueError:
            print("Please enter a valid number.")

def display_qa_menu():
    """Display the question-answer mode."""
    print("\n"+"-"*50)
    print("Question - Answer".center(50))
    print("-"*50)
    print("Type your question or use one of these commands:")
    print("  'help' - Show this help message")
    print("  'back' or 'menu' - Return to main menu")
    print("  'exit' - Exit the program")

def handle_qa_mode(document_text):
    """Handle the question-answering mode with navigation."""
    display_qa_menu()
    
    while True:
        question = input("\nAsk a question: ").strip()
        
        # Handle commands
        if question.lower() in ['back', 'menu']:
            print("\nReturning to menu...")
            return
        elif question.lower() == 'exit':
            print("\nExiting program.. Goodbye!")
            sys.exit(0)
        elif question.lower() == 'help':
            display_qa_menu()
            continue
        elif not question:
            print("Please enter a question or command.")
            continue
        
        # Process question
        print("\nThinking...")
        try:
            answer = answer_question(document_text, question)
            
            print("\nAnswer:")
            # Print the answer with proper wrapping
            for line in textwrap.wrap(str(answer), width=80):
                print(line)
                
            print("\n" + "-"*50)
            print("Type another question, 'back' for menu, or 'help' for options")
            
        except Exception as e:
            print(f"\nError processing question: {e}")
            print("You can try another question or type 'back' to return to the menu")

def main():
    try:
        print("\nWelcome to Reviewer!")
        file_path = get_file_path()
        
        # Extract the text once
        print("\nProcessing file, please wait...")
        document_text = process_file(file_path)
        print("File processed successfully!")
        
        while True:
            choice = display_menu()
            
            if choice == 1:
                print("\nGenerating summary and bullet points, please wait...")
                result = summarize_text_to_bullets(document_text)
                print("\n" + "="*50)
                print("Summary and Key Points".center(50))
                print("="*50)
                print(result)
                
            elif choice == 2:
                handle_qa_mode(document_text)
            
            elif choice == 3:
                print("\nExiting program.. Goodbye!")
                break
    
    except KeyboardInterrupt:
        print("\n\nProgram interrupted by user")
        sys.exit(0)
    except Exception as e:
        print(f"\nAn error occurred: {e}")
        sys.exit(1)

if __name__ == "__main__":
    try:
        main()
    except KeyboardInterrupt:
        print("\nProgram interrupted by user during initialization")
        sys.exit(0)
