import os
from PyPDF2 import PdfReader
from pptx import Presentation
import io
import pytesseract
from pdf2image import convert_from_path
from PIL import Image

from utils.converters import convert_ppt_to_pptx

def process_file(file_path):
    """Extract text based on file type."""
    file_ext = os.path.splitext(file_path)[1].lower()
    if file_ext == '.pdf':
        return extract_text_from_pdf(file_path)
    elif file_ext in ['.pptx', '.ppt']:
        return extract_text_from_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file format: {file_ext}")

def extract_text_from_pdf(file_path):
    """Extract text from each page of the PDF."""
    pdf_text = []
    try:
        with open(file_path, "rb") as file:
            reader = PdfReader(file)
            for page in reader.pages:
                text = page.extract_text()
                if text:  
                    pdf_text.append(text)
    except Exception as e:
        raise IOError(f"Error reading the PDF file: {e}")
    
    # extract text from images in PDF using OCR
    images_text = extract_images_text_from_pdf(file_path)
    if images_text:
        pdf_text.append(images_text)
        
    return " ".join(pdf_text)  

def extract_images_text_from_pdf(file_path):
    """Extract text from images in PDF using OCR."""
    images_text = []
    try:
        # convert pdf pages to images 
        pages = convert_from_path(file_path)
        for i, page in enumerate(pages):
            # extract text from pages in images using ocr
            text = pytesseract.image_to_string(page)
            if text.strip():
                images_text.append(f"[Image text from page {i+1}]: {text}")
    except Exception as e:
        print(f"Warning: Could not extract text from images in PDF: {e}")
    return "\n\n".join(images_text)

def extract_images_text_from_pptx(prs):
    """Extract images from each page of the PowerPoint."""
    images_text = []
    try:
        for i, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    try:
                        # get image data
                        image_data = shape.image.blob
                        image = Image.open(io.BytesIO(image_data))
                        # extract text with ocr 
                        text = pytesseract.image_to_string(image)
                        if text.strip():
                            images_text.append(f"[Image text from slide {i}]: {text}")
                    except Exception as img_error:
                        print(f"Warning: Could not process image in slide {i}: {img_error}")
                        continue
    except Exception as e:
        print(f"Warning: Could not extract texts from images in presentation: {e}")
    return "\n\n".join(images_text)

def extract_text_from_pptx(file_path):
    """Extract text from PowerPoint"""
    pptx_text = []
    file_ext = os.path.splitext(file_path)[1].lower()
    temp_file = None
    try:
        if file_ext == '.ppt':
            print("\nNote: .ppt file detected, attempting automatic conversion...")
            temp_file = convert_ppt_to_pptx(file_path)
            file_path = temp_file
        print("\nExtracting text from presentation...")
        prs = Presentation(file_path)
        
        # process slides
        for slide_number, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                # check if shape has text attribute
                if hasattr(shape, 'text'):
                    try:
                        text = shape.text.strip()
                        if text:
                            slide_text.append(text)
                    except:
                        continue
            if slide_text:
                pptx_text.append(f"[Slide {slide_number}] {' '.join(slide_text)}")
        
        # extract text from images in PPTX using OCR
        images_text = extract_images_text_from_pptx(prs)
        if images_text:
            pptx_text.append(images_text)
            
    except Exception as e:
        error_msg = "Error reading PowerPoint file"
        if file_ext == '.ppt':
            error_msg = "Error reading .ppt file"
        raise IOError(f"{error_msg}: {e}")
    finally:
        # remove temp file if it exists
        if temp_file and os.path.exists(temp_file):
            try:
                os.remove(temp_file)
            except:
                pass
    return "\n\n".join(pptx_text)
